from typing import IO, TextIO
from pptx import Presentation
from datetime import datetime


def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text


def get_words_form_file(file: IO[bytes] | TextIO) -> list[str]:
    print(f"Получение слов из файла {file.name}")

    prs = Presentation(file)

    result: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if is_not_valid(shape.text):
                continue
            result.append(shape.text.strip())

    print(f"Количество : {len(result)} слов")

    return result

def saveTofile(words: list[str]or set[str], filename: str)-> None:
    with open(filename, 'w',encoding='utf-8') as file:
        current_data = datetime.now().data()
        file.writelines(f"{words} : {current_data}"  for w in words)
        print(f"Сохраняем слова: {len(words)} в файл {filename} ")
        print(f"OK")