# Задача 2024.10.23.04
#
# Сохраните из этой презентации только слова, всякие "Правила", "Тур 1" и другие объяснение не нужны.
#
# Сохраните все слова в файл words.txt
# Закомитьте его тоже.

# Задача 2024.10.30.01
#
# После выполнения задачи выше, сохраните полученный список слов в файл words.txt.

# Задача 2024.10.30.02
#
# Напишите тест, который проверяет, что в строке находится всего одно слово.
import zipfile
from zipfile import ZipFile
from pptx import Presentation
from pathlib import Path
import os


def extract():
    zip_name = "croco-blitz-source.zip"
    my_file = ("words.txt")
    prs = Presentation('Osennyaya_igra_3.pptx')

    word = []  #инициализация пустого листа
    folder = os.path.dirname(os.path.realpath(__file__))
    try:
        with zipfile.ZipFile(Path(folder) / 'src' / zip_name) as archive:
            file = archive.namelist()

            presentation_files = [f for f in file if f.endswith('.pptx')]

            if not presentation_files:
                print(f"Файлы презентации не найдены в {zip}")

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text
                        if " " in text or ":" in text or "БЛИЦ-КРОКОДИЛ" in text:
                            continue
                        word.append(text)
                        print(text)

            #
            # for slide in prs.slides:
            #     for shape in slide.shapes:
            #         if not shape.has_text_frame:
            #             continue
            #         print(shape.text)

            with open(my_file, 'w', encoding='utf-8') as f:
                for i in word:
                    f.write(i + '\n')




    except zipfile.BadZipFile as archive:
        print("error")


# def one_word(s):
#     file = open("words.txt")
#     text = file.read()
#     print(text)


if __name__ == '__main__':
    extract()
