from pptx import Presentation


def extract(w: []):
    words = ""
    prs = Presentation(w)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text
                if " " in text or ":" in text or "БЛИЦ-КРОКОДИЛ" in text or "СУПЕРКРОКО" in text:
                    continue
                words += text + '\n'

    return words


