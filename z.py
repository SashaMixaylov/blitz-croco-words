import zipfile
from pptx import Presentation
from pathlib import Path
import os

zip_name = "croco-blitz-source.zip"


def extract_words():
    folder = os.path.dirname(os.path.realpath(__file__))

    try:
        with zipfile.ZipFile(Path(folder) / 'src' / zip_name) as archive:
            file = archive.namelist()

            presentation_files = [f for f in file if f.endswith('.pptx')]

            if not presentation_files:
                print(f"Файлы презентации не найдены в {zip_name}")
                return

            # Извлекаем первую презентацию
            extracted_file = archive.extractall()

            # Создаем экземпляр Presentation
            prs = Presentation('Osennyaya_igra_3.pptx')

            # Извлекаем текст из презентации
            words = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        words.extend(shape.text.split())

            # Фильтруем слова (оставляем только буквы и цифры)
            filtered_words = [''.join(filter(str.isalnum, word)) for word in words]

            # Удаляем пустые строки
            filtered_words = [word for word in filtered_words if word]

            # Записываем результат в файл
            with open('words.txt', 'w', encoding='utf-8') as f:
                for word in filtered_words:
                    f.write(word + '\n')

            print("Извлечение завершено. Слова сохранены в файл words.txt.")

    except Exception as e:
        print(f"Ошибка при извлечении слов: {e}")


if __name__ == '__main__':
    extract_words()
