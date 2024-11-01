# Задача 2024.10.30.04
#
# Теперь попробуйте достать все презентации прямо из архива.
#
# То есть ваша программа должна работать, когда вы поулчаете на вход zip архив, достаёте из него все pptx файлы и парсите каждый,
# сохраняете из него все слова, и после парсинга всех презентаций сохраняете все слова в words.txt

# Задача 2024.11.01.02
#
# В какой контейнер данных нужно собирать слова чтобы исключить дубликаты?
# Поменяйте код так, чтобы это сделать.
#
# Подумайте на каком уровне лучше это делать и почему.
#
# Напишите сколько было слов всего из всех файлов, когда вы сохраняли их в список, и сколько стало, когда вы исключили дубликаты.
import zipfile
from enum import unique

from pptx import Presentation
from pathlib import Path
import os


def extract(zip_path):
    try:
        with zipfile.ZipFile(Path(os.path.dirname(os.path.realpath(__file__))) / 'src' / zip_path) as archive:
            file_list = archive.namelist()
            presentation_files = [file for file in file_list]

            if not presentation_files:
                print(f"Файлы презентации не найдены в {zip_path}")
                return None


            unique_words = set()
            words = []
            for i, presentation_files in enumerate(presentation_files):
                prs = Presentation(archive.open(presentation_files))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                         text = shape.text
                         if " " in text or ":" in text or "БЛИЦ-КРОКОДИЛ" in text or "СУПЕРКРОКО" in text:
                             continue
                         word = text.strip()
                         if word and word not in unique_words:
                             unique_words.add(word)
                         words.append(word)
                         print(word)

            return list(words)
            #ищем дубликаты
            # unique_words = []
            # for w in words:
            #     if word not in unique_words:
            #         unique_words.append(word)
            #         print(f"Дубликаты {w}")


    except zipfile.BadZipFile:
        print("Ошибка при открытии ZIP-файла")
        return None



if __name__ == '__main__':
    zip_name = "croco-blitz-source.zip"
    result = extract(zip_name)

    if result:
        my_file = "words1.txt"
        with open(my_file, 'w', encoding='utf-8') as f:
            for word in result:
                f.write(word + '\n')

        print(f"Извлечение завершено. Слова сохранены в файл {my_file}.")


