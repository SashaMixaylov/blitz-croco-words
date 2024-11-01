# Задача 2024.10.23.06
#
# Разархивируйте архив и достаньте один файл презентации.
#
# Откройте файл презентации в программе с помощью объекта Presentation библиотеки pptx

# Задача 2024.10.23.07
#
# Дополните этот код, положите презентацию рядом с main.py и попробуйте запустить программу и открыть объект Presentation, в нём shape, а в нём slide и в нём найти где лежит текст. С помощью дебаггера.

# Задача 2024.10.25.03
#
# Поменяйте код так, чтобы он выводил текст каждого шэйпа на экран.
# Задача 2024.10.23.04
#
# Сохраните из этой презентации только слова, всякие "Правила", "Тур 1" и другие объяснение не нужны.
#
# Сохраните все слова в файл words.txt
# Закомитьте его тоже.
import zipfile
from pptx import Presentation
from pathlib import Path
import os

zip_name = "croco-blitz-source.zip"
def extract():
    folder = os.path.dirname(os.path.realpath(__file__))
    try:
        with zipfile.ZipFile(Path(folder) / 'src' / zip_name) as archive:
            file = archive.namelist()

            presentation_files = [f for f in file if f.endswith('.pptx')]

            if not presentation_files:
                print(f"Файлы презентации не найдены в {zip}")
                return None

                # Извлекаем первую презентацию
            extracted_file = archive.extract(presentation_files[0])

            # Создаем экземпляр Presentation
            prs = Presentation('Osennyaya_igra_3.pptx')

            # Обрабатываем содержимое презентации
            # for slide in prs.slides:
            #     for shape in slide.shapes:
            #         if hasattr(shape, 'text'):
            #             print(shape.text)
            #
            # return prs
            words = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    print(shape.text)

            filter = ["Правила", "Тур 1"]

            with open('words.txt', 'w', encoding='utf-8') as f:
                for word in filter:
                    f.write(word + '\n')
            print("Найденные слова")


    except zipfile.BadZipFile as archive:
        print("error")



if __name__ == '__main__':
    extract()




